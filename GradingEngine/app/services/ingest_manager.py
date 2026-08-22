from pathlib import Path

from pypdf import PdfReader

from app.services.rag_service import get_collection

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


SUPPORTED_EXTENSIONS = {".pdf", ".pptx"}


def extract_text_from_pptx(file_path: str) -> list[str]:
    """Parse a PowerPoint deck slide-by-slide; one string per slide."""
    if Presentation is None:
        raise ImportError(
            "The 'python-pptx' library is not installed. Run: pip install python-pptx"
        )

    prs = Presentation(file_path)
    slide_texts: list[str] = []

    for slide in prs.slides:
        slide_content: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        slide_texts.append("\n".join(slide_content))

    return slide_texts


def _index_chunk(
    collection,
    *,
    doc_id: str,
    text: str,
    course_name: str,
    page_number: int,
    source_file: str,
    doc_type: str,
) -> None:
    collection.upsert(
        documents=[text.strip()],
        metadatas=[
            {
                "course": course_name,
                "page_number": page_number,
                "source_file": source_file,
                "type": doc_type,
            }
        ],
        ids=[doc_id],
    )


def process_and_index_lecture(file_path: str, course_name: str) -> int:
    """
    Index a lecture PDF or PPTX into ChromaDB (one vector per page or slide).
    Returns the number of pages/slides indexed.
    """
    path = Path(file_path)
    if not path.is_file():
        raise FileNotFoundError(f"Target file not found at {file_path}")

    extension = path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{extension}'. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    print(f"Parsing document: '{path.name}' ({extension})...")
    collection = get_collection()
    indexed_count = 0
    course_slug = course_name.lower().replace(" ", "_")
    file_slug = "".join(ch if ch.isalnum() else "_" for ch in path.stem).strip("_").lower()[:80] or "file"

    if extension == ".pdf":
        reader = PdfReader(str(path))
        for page_num, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text() or ""
            if not page_text.strip():
                continue

            doc_id = f"{course_slug}_pdf_{file_slug}_p{page_num}"
            _index_chunk(
                collection,
                doc_id=doc_id,
                text=page_text,
                course_name=course_name,
                page_number=page_num,
                source_file=path.name,
                doc_type="PDF",
            )
            indexed_count += 1

    elif extension == ".pptx":
        slide_texts = extract_text_from_pptx(str(path))
        for slide_num, slide_text in enumerate(slide_texts, start=1):
            if not slide_text.strip():
                continue

            doc_id = f"{course_slug}_pptx_{file_slug}_s{slide_num}"
            _index_chunk(
                collection,
                doc_id=doc_id,
                text=slide_text,
                course_name=course_name,
                page_number=slide_num,
                source_file=path.name,
                doc_type="PPTX",
            )
            indexed_count += 1

    print(f"Vector ingestion completed. Registered {indexed_count} items.")
    return indexed_count


# Backward-compatible alias used by existing imports.
def process_and_index_pdf(file_path: str, course_name: str) -> int:
    return process_and_index_lecture(file_path, course_name)
